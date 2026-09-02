import io

from fastapi import HTTPException, UploadFile
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from pathlib import Path


class StudyMaterialService:

    async def extract_text(
        self,
        file: UploadFile,
    ) -> str:

        content = await file.read()

        if not content:
            raise HTTPException(
                status_code=400,
                detail=(
                    f"File is empty: "
                    f"{file.filename}"
                ),
            )

        filename = file.filename or ""

        extension = Path(
            filename
        ).suffix.lower()

        if extension == ".pdf":

            text = self._extract_pdf(
                content
            )

        elif extension == ".txt":

            text = self._extract_txt(
                content
            )

        elif extension == ".docx":

            text = self._extract_docx(
                content
            )

        elif extension == ".pptx":

            text = self._extract_pptx(
                content
            )

        else:

            raise HTTPException(
                status_code=400,
                detail=(
                    f"Unsupported file type: "
                    f"{filename}"
                ),
            )

        text = self._clean_text(
            text
        )

        if not text:
            raise HTTPException(
                status_code=400,
                detail=(
                    f"Could not extract usable text from: "
                    f"{filename}"
                ),
            )

        return text

    def _extract_pdf(
        self,
        content: bytes,
    ) -> str:

        try:
            reader = PdfReader(
                io.BytesIO(content)
            )

            return "\n".join(
                page.extract_text() or ""
                for page in reader.pages
            )

        except Exception as error:
            raise HTTPException(
                status_code=400,
                detail="Failed to read PDF file",
            ) from error

    def _extract_txt(
        self,
        content: bytes,
    ) -> str:

        try:
            return content.decode("utf-8")

        except UnicodeDecodeError:

            try:
                return content.decode(
                    "utf-8-sig"
                )

            except UnicodeDecodeError as error:
                raise HTTPException(
                    status_code=400,
                    detail=(
                        "Text file must use UTF-8 encoding"
                    ),
                ) from error

    def _extract_docx(
        self,
        content: bytes,
    ) -> str:

        try:
            document = Document(
                io.BytesIO(content)
            )

            return "\n".join(
                paragraph.text
                for paragraph in document.paragraphs
                if paragraph.text.strip()
            )

        except Exception as error:
            raise HTTPException(
                status_code=400,
                detail="Failed to read DOCX file",
            ) from error

    def _extract_pptx(
        self,
        content: bytes,
    ) -> str:

        try:
            presentation = Presentation(
                io.BytesIO(content)
            )

            slide_text = []

            for slide in presentation.slides:

                for shape in slide.shapes:

                    if hasattr(
                        shape,
                        "text",
                    ):
                        text = shape.text.strip()

                        if text:
                            slide_text.append(
                                text
                            )

            return "\n".join(
                slide_text
            )

        except Exception as error:
            raise HTTPException(
                status_code=400,
                detail="Failed to read PPTX file",
            ) from error

    def _clean_text(
        self,
        text: str,
    ) -> str:

        lines = [
            line.strip()
            for line in text.splitlines()
            if line.strip()
        ]

        return "\n".join(lines)